import streamlit as st
import os
import tempfile
import whisper
from PyPDF2 import PdfReader
from moviepy.editor import *
import docx2txt
import io
import platform
import subprocess
import time
from datetime import datetime

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from dotenv import load_dotenv
import google.generativeai as genai

from docx import Document
from pptx import Presentation

# ========== Load API Key ==========
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# ========== Utilities ==========
def get_libreoffice_path():
    """Get the correct LibreOffice path based on the operating system"""
    system = platform.system()
    
    if system == "Windows":
        # Common Windows installation paths for LibreOffice
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice*\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice*\program\soffice.exe"
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                return path
            # Try with wildcard expansion
            import glob
            expanded_paths = glob.glob(path)
            if expanded_paths:
                return expanded_paths[0]
    else:
        # For macOS/Linux, just use the command name
        return "soffice"
    
    return None

def convert_doc_to_pdf(doc_path):
    """Convert .doc file to PDF using LibreOffice"""
    # Get directory and filename
    doc_dir = os.path.dirname(doc_path)
    temp_output_dir = os.path.join(doc_dir, "temp_output")
    os.makedirs(temp_output_dir, exist_ok=True)
    
    # Get the correct LibreOffice path
    libreoffice_path = get_libreoffice_path()
    if not libreoffice_path:
        st.error("LibreOffice not found. Please install LibreOffice to process .doc files.")
        return None
    
    try:
        # Prepare the command based on the operating system
        if platform.system() == "Windows":
            # Windows needs the full path and different path separators
            cmd = [
                libreoffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_output_dir.replace("/", "\\"),
                doc_path.replace("/", "\\")
            ]
        else:
            # Unix-like systems (macOS/Linux)
            cmd = [
                libreoffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_output_dir,
                doc_path
            ]
        
        # Run the conversion command
        subprocess.run(cmd, check=True, capture_output=True)
        
        # Get the PDF file path
        base_name = os.path.basename(doc_path)
        pdf_name = os.path.splitext(base_name)[0] + ".pdf"
        pdf_path = os.path.join(temp_output_dir, pdf_name)
        
        # Check if conversion was successful
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception as e:
        st.error(f"LibreOffice conversion failed: {str(e)}")
        if platform.system() == "Windows":
            st.info("If you're having issues with LibreOffice on Windows, try:")
            st.info("1. Make sure LibreOffice is installed in the default location")
            st.info("2. Try running the application as administrator")
            st.info("3. Check if LibreOffice is properly registered in the system PATH")
    
    return None

def extract_text_from_doc(doc):
    """Extract text from a .doc file via PDF conversion"""
    # Create a temporary directory to work with the file
    with tempfile.TemporaryDirectory() as temp_dir:
        # Save the uploaded file to the temp directory
        temp_file_path = os.path.join(temp_dir, "temp.doc")
        with open(temp_file_path, "wb") as f:
            f.write(doc.getvalue())
        
        # Step 1: Try PDF conversion first
        pdf_path = convert_doc_to_pdf(temp_file_path)
        if pdf_path and os.path.exists(pdf_path):
            try:
                # Extract text from the PDF
                pdf_reader = PdfReader(pdf_path)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                
                if text.strip():
                    return text
            except Exception as e:
                st.warning(f"PDF extraction failed: {str(e)}")
        
        # Step 2: Fallback to previous methods if PDF conversion fails
        
        # Try using docx2txt even though it's not ideal for .doc
        try:
            text = docx2txt.process(temp_file_path)
            if text.strip():  # If we got meaningful text
                return text
        except Exception:
            pass
        
        # Try catdoc if available (not on Windows)
        if platform.system() != "Windows":
            try:
                result = subprocess.run(
                    ["catdoc", temp_file_path], 
                    capture_output=True, 
                    text=True, 
                    check=True
                )
                if result.stdout.strip():
                    return result.stdout
            except Exception:
                pass
        
        # Try to convert directly to TXT using LibreOffice
        try:
            txt_output_dir = os.path.join(temp_dir, "txt_output")
            os.makedirs(txt_output_dir, exist_ok=True)
            
            libreoffice_path = get_libreoffice_path()
            if libreoffice_path:
                if platform.system() == "Windows":
                    cmd = [
                        libreoffice_path,
                        "--headless",
                        "--convert-to", "txt:Text",
                        "--outdir", txt_output_dir.replace("/", "\\"),
                        temp_file_path.replace("/", "\\")
                    ]
                else:
                    cmd = [
                        libreoffice_path,
                        "--headless",
                        "--convert-to", "txt:Text",
                        "--outdir", txt_output_dir,
                        temp_file_path
                    ]
                
                subprocess.run(cmd, check=True, capture_output=True)
                
                base_name = os.path.basename(temp_file_path)
                txt_name = os.path.splitext(base_name)[0] + ".txt"
                txt_path = os.path.join(txt_output_dir, txt_name)
                
                if os.path.exists(txt_path):
                    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()
                        if text.strip():
                            return text
        except Exception:
            pass
        
    # If all methods fail
    return ""

def check_doc_conversion_tools():
    """Check if necessary tools for .doc conversion are installed"""
    tools_status = {
        "libreoffice": False,
        "catdoc": False
    }
    
    # Check for LibreOffice
    libreoffice_path = get_libreoffice_path()
    if libreoffice_path:
        try:
            if platform.system() == "Windows":
                # On Windows, we need to use the full path
                subprocess.run([libreoffice_path, "--version"], capture_output=True, check=False)
            else:
                subprocess.run(["soffice", "--version"], capture_output=True, check=False)
            tools_status["libreoffice"] = True
        except Exception:
            pass
    
    # Check for catdoc (only on non-Windows systems)
    if platform.system() != "Windows":
        try:
            subprocess.run(["catdoc", "--version"], capture_output=True, check=False)
            tools_status["catdoc"] = True
        except FileNotFoundError:
            pass
    
    return tools_status

# ========== File Text Extraction ==========
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        try:
            pdf_reader = PdfReader(pdf)
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            st.success(f"Successfully processed PDF: {pdf.name}")
        except Exception as e:
            st.error(f"Error processing PDF {pdf.name}: {str(e)}")
    
    if not text.strip():
        raise ValueError("No extractable text found in PDF files.")
    return text

def get_docx_text(docx_file):
    try:
        document = Document(docx_file)
        text = "\n".join([para.text for para in document.paragraphs])
        st.success(f"Successfully processed DOCX: {docx_file.name}")
        return text
    except Exception as e:
        st.error(f"Error processing DOCX file {docx_file.name}: {str(e)}")
        return ""

def get_pptx_text(pptx_file):
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        st.success(f"Successfully processed PPTX: {pptx_file.name}")
        return text
    except Exception as e:
        st.error(f"Error processing PPTX file {pptx_file.name}: {str(e)}")
        return ""

def get_txt_text(txt_file):
    try:
        text = txt_file.read().decode("utf-8")
        st.success(f"Successfully processed TXT: {txt_file.name}")
        return text
    except Exception as e:
        st.error(f"Error processing TXT file {txt_file.name}: {str(e)}")
        return ""

def get_doc_text(doc_file):
    """Extract text from a .doc file"""
    st.info(f"Processing .DOC file: {doc_file.name} - this may take a moment...")
    doc_text = extract_text_from_doc(doc_file)
    
    if doc_text and len(doc_text.strip()) > 0:
        st.success(f"Successfully processed DOC: {doc_file.name}")
        return doc_text
    else:
        st.error(f"Could not extract text from {doc_file.name}")
        st.warning("Try converting your .doc file to .docx or .pdf format before uploading")
        return ""

# ========== Transcription ==========
def transcribe_video(video_file):
    st.info(f"Transcribing video: {video_file.name} - this may take a moment...")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_video:
        temp_video.write(video_file.read())
        temp_path = temp_video.name

    model = whisper.load_model("base")
    result = model.transcribe(temp_path)

    if "segments" not in result or not result["segments"]:
        st.error(f"No speech content found in video: {video_file.name}")
        raise ValueError("No speech content found in the video.")

    transcript = ""
    metadata = []
    for segment in result["segments"]:
        start_time = segment["start"]
        text = segment["text"].strip()
        if text:
            transcript += f"{text} "
            metadata.append({"start_time": start_time})

    st.success(f"Successfully transcribed video: {video_file.name}")
    return transcript.strip(), metadata

def transcribe_audio(audio_file):
    st.info(f"Transcribing audio: {audio_file.name} - this may take a moment...")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_audio:
        temp_audio.write(audio_file.read())
        temp_path = temp_audio.name

    model = whisper.load_model("base")
    result = model.transcribe(temp_path)

    if "segments" not in result or not result["segments"]:
        st.error(f"No speech content found in audio: {audio_file.name}")
        raise ValueError("No speech content found in the audio.")

    transcript = ""
    metadata = []
    for segment in result["segments"]:
        start_time = segment["start"]
        text = segment["text"].strip()
        if text:
            transcript += f"{text} "
            metadata.append({"start_time": start_time})

    st.success(f"Successfully transcribed audio: {audio_file.name}")
    return transcript.strip(), metadata

# ========== Chunking, Embedding, QA ==========
def get_text_chunks(text, chunk_size=1000, chunk_overlap=100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

def get_vector_store(text_chunks, metadatas=None, index_name="faiss_index"):
    if not text_chunks:
        raise ValueError("No text chunks to index!")

    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    if metadatas and len(metadatas) != len(text_chunks):
        raise ValueError("Length of metadata must match text chunks.")

    vector_store = FAISS.from_texts(
        text_chunks,
        embedding=embeddings,
        metadatas=metadatas if metadatas else None
    )
    vector_store.save_local(index_name)

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context.
    If the answer is not available in the context, say "Answer is not available in the context."
    
    Context:\n{context}\n
    Question:\n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="models/gemini-1.5-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def answer_question(user_question, index_name="faiss_index"):
    # Remove any progress indicators before showing the answer
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    # Remove any spinners or progress indicators
    vector_store = FAISS.load_local(index_name, embeddings, allow_dangerous_deserialization=True)
    docs = vector_store.similarity_search(user_question)

    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)

    # st.markdown('<div class="answer-box">', unsafe_allow_html=True)
    st.markdown("### 📜 Reply:")
    st.write(response["output_text"])
    st.markdown('</div>', unsafe_allow_html=True)

    # for doc in docs:
    #     metadata = doc.metadata
    #     if "start_time" in metadata:
    #         mins = int(metadata["start_time"] // 60)
    #         secs = int(metadata["start_time"] % 60)
    #         st.markdown(f'<div style="color:#6b7280; font-size:0.9rem;">⏱️ Reference timestamp: {mins:02d}:{secs:02d}</div>', unsafe_allow_html=True)
    #         break

# ========== Streamlit App ==========
def main():
    # Set page configuration
    st.set_page_config(
        page_title="Mediaverse",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Custom CSS for a professional look
    st.markdown("""
    <style>
    /* Main app styling */
    .main {
        background-color: #f8fafc;
    }
    .main .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
        max-width: 1200px;
    }
    
    /* File uploader styling - customize the buttons */
    /* Target all file uploader buttons */
    .stFileUploader button,
    .stFileUploader [role="button"],
    button.css-1k0ckh2,
    [data-testid="stFileUploader"] button,
    [data-testid="stFileUploadDropzone"] button,
    div[data-testid="stFileUploader"] button {
        background-color: #e0f2fe !important;
        color: #5b5b5b !important;
        border: 1px solid #4CB4FA !important;
        border-radius: 6px !important;
        padding: 0.5rem 1rem !important;
        font-weight: 500 !important;
        transition: all 0.3s ease !important;
    }
    
    /* File uploader button hover effect - use multiple selectors to ensure it works */
    .stFileUploader button:hover,
    .stFileUploader [role="button"]:hover,
    button.css-1k0ckh2:hover,
    [data-testid="stFileUploader"] button:hover,
    [data-testid="stFileUploadDropzone"] button:hover,
    div[data-testid="stFileUploader"] button:hover {
        background-color: transparent !important;
        color: #5b5b5b !important;
        border: 1px solid #4CB4FA !important;
    }
    
    /* Additional selector for the 'Browse files' button */
    button.st-emotion-cache-1q5lpud {
        background-color: #b5a0ca !important;
        color: white !important;
        border: 1px solid #b5a0ca !important;
    }
    
    button.st-emotion-cache-1q5lpud:hover {
        background-color: transparent !important;
        color: #b5a0ca !important;
        border: 1px solid #b5a0ca !important;
    }
    
    /* File uploader area styling */
    .stFileUploader > div:nth-child(3) {
        background-color: #f8f9fa !important;
        border: 2px dashed #b5a0ca !important;
        border-radius: 10px !important;
        padding: 20px !important;
    }
    
    /* Add gap between uploaded file and upload section */
    [data-testid="stFileUploader"] > div:last-child {
        margin-top: 16px !important; 
    }
    
    /* Style for uploaded file name container - reduce padding */
    .uploadedFile {
        padding: 6px 10px !important;
        margin: 4px 0 !important;
        background-color: #f8fafc !important;
        border: 1px solid #e2e8f0 !important;
        border-radius: 6px !important;
    }
    
    /* Make file name text smaller */
    .uploadedFileName {
        font-size: 0.9rem !important;
        padding: 2px !important;
    }
    
    /* Reduce padding around horizontal lines */
    hr {
        margin: 0.5rem 0 !important;
        padding: 0 !important;
        border-color: #d1c4ce !important;
        opacity: 0.9 !important;
    }
    
    /* Reduce padding for markdown horizontal rules */
    .element-container:has(hr) {
        margin-top: 0.5rem !important;
        margin-bottom: 0.5rem !important;
        padding-top: 0 !important;
        padding-bottom: 0 !important;
    }
    
    /* Main content area */
    .content-container {
        background-color: white;
        border-radius: 10px;
        padding: 2rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);
        margin-bottom: 2rem;
    }
    
    /* Headers and text */
    h1, h2, h3 {
        font-family: 'Helvetica Neue', sans-serif;
    }
    
    h1 {
        color: #1e40af;
        font-weight: 700;
        margin-bottom: 1.5rem !important;
    }
    
    .subheader {
        color: #64748b;
        font-size: 1.2rem;
        margin-bottom: 2rem;
    }
    
    /* Sidebar styling */
    section[data-testid="stSidebar"] {
        background: linear-gradient(90deg, #e0f2fe 0%, #f0f9ff 50%, #f1f5f9 100%);
        border-right: 1px solid #e2e8f0;
    }
    
    section[data-testid="stSidebar"] > div {
        padding-top: 2rem;
        padding-left: 1.5rem;
        padding-right: 1.5rem;
    }
    
    .sidebar-header {
        font-size: 1.2rem;
        font-weight: 600;
        color: #1e40af;
        margin-top: 1rem;
        margin-bottom: 1rem;
        padding-bottom: 0.5rem;
        border-bottom: 1px solid #e2e8f0;
    }
    
    .sidebar-subheader {
        font-size: 1rem;
        font-weight: 600;
        color: #475569;
        margin-top: 1.5rem;
        margin-bottom: 0.5rem;
        display: flex;
        align-items: center;
    }
    
    /* Custom Submit and Process button */
    div[data-testid="stButton"] button[kind="secondary"] {
        background-color: white;
       color: #5b5b5b !important;
        border: 1px solid #4CB4FA !important;
        transition: background-color 0.3s;
    }
    
    div[data-testid="stButton"] button[kind="secondary"]:hover {
        background-color: #e0f2fe;
    }

    /* File status containers */
    .status-container {
        background-color: white;
        border-radius: 8px;
        padding: 0.75rem;
        margin-bottom: 1rem;
        border: 1px solid #e2e8f0;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.05);
    }
    
    /* Success message styling */
    .success-msg {
        color: #047857;
        font-weight: 500;
        display: flex;
        align-items: center;
        margin: 0.5rem 0;
    }
    
    /* Text input styling */
    .stTextInput input {
        background: linear-gradient(90deg, #e0f2fe 0%, #f0f9ff 50%, #f1f5f9 100%);
    }

    /* Button styling */
    .stButton>button {
        background-color: #1e40af;
        color: white;
        font-weight: 500;
        border: none;
        border-radius: 6px;
        padding: 0.6rem 1rem;
        width: 100%;
        transition: all 0.3s ease;
    }
    
    .stButton>button:hover {
        background-color: #1e3a8a;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        transform: translateY(-2px);
    }
    
    /* Input box styling */
    .stTextInput>div>div>input {
        border-radius: 6px;
        border: 1px solid #cbd5e1;
        padding: 0.75rem !important;
        box-shadow: 0 1px 2px rgba(0, 0, 0, 0.05);
    }
    
    .stTextInput>div>div>input:focus {
        border-color: #3b82f6;
        box-shadow: 0 0 0 3px rgba(59, 130, 246, 0.2);
    }
    
    /* Footer styling */
    .footer {
        position: fixed;
        bottom: 0;
        left: 0;
        width: 100%;
        background: linear-gradient(90deg, #e0f2fe 0%, #f0f9ff 50%, #f1f5f9 100%);
        padding: 0.75rem;
        text-align: center;
        font-size: 0.8rem;
        color: #64748b;
        border-top: 1px solid #e2e8f0;
        z-index: 1000;
    }
    
    /* Progress bar styling - custom green color */
    .stProgress > div > div > div > div {
        background-color: #10b981 !important;
    }
    
    /* Progress bar container */
    .stProgress > div > div > div {
        background-color: #e2e8f0;
        height: 8px !important;
        border-radius: 4px !important;
    }
    
    /* Make progress bars visible */
    .stProgress {
        display: block !important;
    }
    
    /* Answer box */
    .answer-box {
        background-color: #f0f9ff;
        border-left: 4px solid #0ea5e9;
        padding: 1.25rem;
        border-radius: 0 8px 8px 0;
        margin: 1rem 0;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.05);
    }
    
    /* Hide any automatic spinners */
    .stSpinner {
        display: none !important;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # App header with icon and gradient background
    st.markdown("""
    <div style="
        padding: 1.5rem; 
        border-radius: 10px; 
        margin-bottom: 2rem;
        background: linear-gradient(90deg, #e0f2fe 0%, #f0f9ff 50%, #f1f5f9 100%);
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);
        display: flex;
        align-items: center;
    ">
        <div style="font-size:3.5rem; margin-right: 1.5rem;">📚</div>
        <div>
            <h1 style="margin: 0; padding: 0;"> Mediaverse </h1>
            <p style="color: #64748b; font-size: 1.2rem; margin-top: 0.5rem;">
                Upload documents, videos, or audio files and chat with their content using AI
            </p>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Create a container for the main content
    with st.container():
        # st.markdown('<div class="content-container">', unsafe_allow_html=True)
        user_question = st.text_input("💬 Ask a question about your documents:", 
                                    placeholder="e.g., 'What are the key points in the document?'",
                                    help="Type your question here and press Enter")

        if user_question:
            question_start_time = time.time()
            st.session_state.question_time = question_start_time
            try:
                # Directly call answer_question without any containers or spinners
                answer_question(user_question)
                
                # Calculate and display question-answer duration
                if st.session_state.question_time:
                    answer_time = time.time()
                    qa_duration = answer_time - st.session_state.question_time
                    st.info(f"⏳ Response time: {qa_duration:.2f} seconds")
                    # Reset question time for next question
                    st.session_state.question_time = None
            except Exception as e:
                st.error(f"Error during question answering: {e}")
        st.markdown('</div>', unsafe_allow_html=True)
    
    # Sidebar for file uploads
    with st.sidebar:
        st.markdown('<div class="sidebar-header">📁 Document Upload Center</div>', unsafe_allow_html=True)
        st.markdown("Upload your files below and click the Submit and Process button")
        
        # PDF section
        st.markdown('<div class="sidebar-subheader">📄 PDF Documents</div>', unsafe_allow_html=True)
        pdf_docs = st.file_uploader("Upload PDF files", type=["pdf"], accept_multiple_files=True)
        pdf_status = st.empty()
        st.markdown("---")
        
        # DOC section - Updated to accept multiple files
        st.markdown('<div class="sidebar-subheader">📝 Word Documents</div>', unsafe_allow_html=True)
        doc_files = st.file_uploader("Upload DOC/DOCX files", type=["docx", "doc"], accept_multiple_files=True)
        doc_status = st.empty()
        st.markdown("---")

        # PPT section - Updated to accept multiple files
        st.markdown('<div class="sidebar-subheader">🖼️ Presentations</div>', unsafe_allow_html=True)
        ppt_files = st.file_uploader("Upload PPT/PPTX files", type=["pptx", "ppt"], accept_multiple_files=True)
        ppt_status = st.empty()
        st.markdown("---")
        
        # TXT section
        st.markdown('<div class="sidebar-subheader">📋 Text Files</div>', unsafe_allow_html=True)
        txt_file = st.file_uploader("Upload TXT file", type=["txt"], accept_multiple_files=True)
        txt_status = st.empty()
        st.markdown("---")
        
        # Video section
        st.markdown('<div class="sidebar-subheader">🎬 Video Files</div>', unsafe_allow_html=True)
        video_files = st.file_uploader("Upload video files", type=["mp4", "mov", "avi"], accept_multiple_files=True)
        video_status = st.empty()
        st.markdown("---")
        
        # Audio section
        st.markdown('<div class="sidebar-subheader">🎵 Audio Files</div>', unsafe_allow_html=True)
        audio_files = st.file_uploader("Upload audio files", type=["mp3", "wav", "m4a"], accept_multiple_files=True)
        audio_status = st.empty()
        st.markdown("---")
        
        # Add some space before the button
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Process button with gradient styling
        st.markdown("""
        <style>
        div.stButton > button:first-child {
            background: linear-gradient(90deg, #1e40af 0%, #3b82f6 100%);
            color: white;
            font-weight: 600;
            padding: 0.6rem 1rem;
            border-radius: 6px;
            border: none;
            width: 100%;
            transition: all 0.3s ease;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        }
        div.stButton > button:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 8px rgba(0, 0, 0, 0.15);
        }
        </style>
        """, unsafe_allow_html=True)
        
        # Initialize session state for timestamps if not exists
        if 'question_time' not in st.session_state:
            st.session_state.question_time = None
        if 'process_start_time' not in st.session_state:
            st.session_state.process_start_time = None

        # Process uploaded files
        if st.button("🚀 Submit and Process", help="Click to process all uploaded files"):
            # Record start time for processing
            st.session_state.process_start_time = time.time()
            try:
                all_chunks = []
                all_metadata = []
                files_processed = False
                
                # Process PDF files
                if pdf_docs:
                    with pdf_status.container():
                        # Create a placeholder for the processing message
                        pdf_msg = st.empty()
                        pdf_msg.write("⏳ Processing PDF files...")
                        
                        # Add progress bar for PDF processing
                        pdf_progress = st.progress(0)
                        pdf_text = ""
                        
                        total_pdfs = len(pdf_docs)
                        
                        for i, pdf in enumerate(pdf_docs):
                            try:
                                # Update progress in steps of 20%, 40%, 60%, 80%, 100%
                                progress_percent = min(100, int(((i+1) / total_pdfs) * 100))
                                # Round to nearest 20%
                                stepped_progress = min(100, ((progress_percent + 10) // 20) * 20)
                                pdf_progress.progress(stepped_progress)
                                
                                pdf_reader = PdfReader(pdf)
                                pdf_content = ""
                                for page in pdf_reader.pages:
                                    pdf_content += page.extract_text() or ""
                                if pdf_content.strip():
                                    pdf_text += pdf_content
                                    st.markdown(f'<div class="success-msg">✅ Successfully processed: {pdf.name}</div>', unsafe_allow_html=True)
                                else:
                                    st.warning(f"No text content found in PDF: {pdf.name}")
                            except Exception as e:
                                st.error(f"Error processing PDF {pdf.name}: {str(e)}")
                        
                        # Ensure the progress bar is at 100% when done
                        pdf_progress.progress(100)
                        pdf_msg.empty()  # Clear the processing message
                        
                        if pdf_text.strip():
                            pdf_chunks = get_text_chunks(pdf_text)
                            all_chunks.extend(pdf_chunks)
                            all_metadata.extend([{}] * len(pdf_chunks))
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the PDF files.")

                # Process DOC/DOCX files
                if doc_files:
                    with doc_status.container():
                        # Create a placeholder for the processing message
                        doc_msg = st.empty()
                        doc_msg.write("⏳ Processing document files...")
                        
                        # Add progress bar
                        doc_progress = st.progress(0)
                        doc_text = ""
                        
                        total_docs = len(doc_files)
                        
                        # First step - 20%
                        doc_progress.progress(20)
                        
                        for i, doc_file in enumerate(doc_files):
                            # Calculate progress percentage
                            progress_percent = min(100, int(((i+1) / total_docs) * 80) + 20)  # Start at 20%, go to 100%
                            # Round to nearest 20%
                            stepped_progress = min(100, ((progress_percent + 10) // 20) * 20)
                            doc_progress.progress(stepped_progress)
                            
                            if doc_file.name.lower().endswith('.docx'):
                                try:
                                    # Process DOCX file
                                    document = Document(doc_file)
                                    file_text = "\n".join([para.text for para in document.paragraphs])
                                    if file_text.strip():
                                        doc_text += file_text + "\n\n"
                                        st.markdown(f'<div class="success-msg">✅ Successfully processed: {doc_file.name}</div>', unsafe_allow_html=True)
                                    else:
                                        st.warning(f"No text content found in DOCX: {doc_file.name}")
                                except Exception as e:
                                    st.error(f"Error processing DOCX file {doc_file.name}: {str(e)}")
                            else:  # .doc file
                                # Create a placeholder for doc processing message
                                doc_file_msg = st.empty()
                                doc_file_msg.info(f"Processing .DOC file: {doc_file.name} - this may take a moment...")
                                try:
                                    file_text = extract_text_from_doc(doc_file)
                                    # Clear the processing message for this specific file
                                    doc_file_msg.empty()
                                    if file_text and len(file_text.strip()) > 0:
                                        doc_text += file_text + "\n\n"
                                        st.markdown(f'<div class="success-msg">✅ Successfully processed: {doc_file.name}</div>', unsafe_allow_html=True)
                                    else:
                                        st.error(f"Could not extract text from {doc_file.name}")
                                        st.warning("Try converting your .doc file to .docx or .pdf format before uploading")
                                except Exception as e:
                                    # Clear the processing message in case of error
                                    doc_file_msg.empty()
                                    st.error(f"Error processing DOC file {doc_file.name}: {str(e)}")
                        
                        # Complete the progress bar and clear the main processing message
                        doc_progress.progress(100)
                        doc_msg.empty()  # Clear the processing message
                        
                        if doc_text.strip():  # Only process if we got some text
                            doc_chunks = get_text_chunks(doc_text)
                            all_chunks.extend(doc_chunks)
                            all_metadata.extend([{}] * len(doc_chunks))
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the document files.")

                # Process PPT/PPTX files
                if ppt_files:
                    with ppt_status.container():
                        # Create a placeholder for the processing message
                        ppt_msg = st.empty()
                        ppt_msg.write("⏳ Processing presentation files...")
                        
                        # Add progress bar
                        ppt_progress = st.progress(0)
                        ppt_text = ""
                        
                        total_ppts = len(ppt_files)
                        
                        # First step - 20%
                        ppt_progress.progress(20)
                        
                        for i, ppt_file in enumerate(ppt_files):
                            # Update progress in steps
                            progress_percent = min(100, int(((i+1) / total_ppts) * 80) + 20)  # Start at 20%, go to 100%
                            # Round to nearest 20%
                            stepped_progress = min(100, ((progress_percent + 10) // 20) * 20)
                            ppt_progress.progress(stepped_progress)
                            
                            try:
                                prs = Presentation(ppt_file)
                                file_text = ""
                                
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            file_text += shape.text + "\n"
                                
                                if file_text.strip():
                                    ppt_text += file_text + "\n\n"
                                    st.markdown(f'<div class="success-msg">✅ Successfully processed: {ppt_file.name}</div>', unsafe_allow_html=True)
                                else:
                                    st.warning(f"No text content found in PPT/PPTX: {ppt_file.name}")
                            except Exception as e:
                                st.error(f"Error processing PPT/PPTX file {ppt_file.name}: {str(e)}")
                        
                        # Complete the progress bar and clear the processing message
                        ppt_progress.progress(100)
                        ppt_msg.empty()  # Clear the processing message
                        
                        if ppt_text.strip():
                            ppt_chunks = get_text_chunks(ppt_text)
                            all_chunks.extend(ppt_chunks)
                            all_metadata.extend([{}] * len(ppt_chunks))
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the presentation files.")

                # Process TXT files
                if txt_file:
                    with txt_status.container():
                        # Create a placeholder for the processing message
                        txt_msg = st.empty()
                        txt_msg.write("⏳ Processing text files...")
                        
                        # Add progress bar with initial 20%
                        txt_progress = st.progress(20)
                        txt_text = ""
                        
                        # If txt_file is a list (multiple files)
                        if isinstance(txt_file, list):
                            total_txts = len(txt_file)
                            
                            for i, txt in enumerate(txt_file):
                                # Update progress in steps
                                progress_percent = min(100, int(((i+1) / total_txts) * 80) + 20)  # Start at 20%, go to 100%
                                # Round to nearest 20%
                                stepped_progress = min(100, ((progress_percent + 10) // 20) * 20)
                                txt_progress.progress(stepped_progress)
                                
                                try:
                                    file_text = txt.read().decode("utf-8")
                                    if file_text.strip():
                                        txt_text += file_text + "\n\n"
                                        st.markdown(f'<div class="success-msg">✅ Successfully processed: {txt.name}</div>', unsafe_allow_html=True)
                                    else:
                                        st.warning(f"TXT file is empty: {txt.name}")
                                except Exception as e:
                                    st.error(f"Error processing TXT file {txt.name}: {str(e)}")
                        else:
                            # Single file processing - show stepped progress
                            txt_progress.progress(40)  # Read file - 40%
                            try:
                                txt_text = txt_file.read().decode("utf-8")
                                txt_progress.progress(80)  # Complete reading - 80%
                                if txt_text.strip():
                                    st.markdown(f'<div class="success-msg">✅ Successfully processed: {txt_file.name}</div>', unsafe_allow_html=True)
                                else:
                                    st.warning(f"TXT file is empty: {txt_file.name}")
                            except Exception as e:
                                st.error(f"Error processing TXT file {txt_file.name}: {str(e)}")
                        
                        # Complete the progress bar and clear the processing message
                        txt_progress.progress(100)  # Complete - 100%
                        txt_msg.empty()  # Clear the processing message
                        
                        if txt_text.strip():
                            txt_chunks = get_text_chunks(txt_text)
                            all_chunks.extend(txt_chunks)
                            all_metadata.extend([{}] * len(txt_chunks))
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the text files.")

                # Process video files
                if video_files:
                    with video_status.container():
                        # Create a placeholder for the processing message
                        video_msg = st.empty()
                        video_msg.write("⏳ Processing video files...")
                        
                        # Add progress bar
                        video_progress = st.progress(0)
                        all_transcripts = ""
                        all_video_metadata = []
                        
                        if isinstance(video_files, list):
                            total_videos = len(video_files)
                            
                            for i, video_file in enumerate(video_files):
                                # Calculate step progress for this file
                                file_start_progress = min(100, int((i / total_videos) * 100))
                                file_end_progress = min(100, int(((i+1) / total_videos) * 100))
                                
                                # Round to nearest 20% for display
                                stepped_start = ((file_start_progress + 10) // 20) * 20
                                stepped_end = ((file_end_progress + 10) // 20) * 20
                                
                                # Show the current step progress
                                video_progress.progress(stepped_start)
                                
                                # Create placeholder for model loading and transcription messages
                                model_msg = st.empty()
                                transcribe_msg = st.empty()
                                
                                try:
                                    # Load model - 20% step for this file
                                    model_msg.info("Loading transcription model...")
                                    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_video:
                                        temp_video.write(video_file.read())
                                        temp_path = temp_video.name
                                    
                                    model = whisper.load_model("base")
                                    model_msg.empty()  # Clear the model loading message
                                    
                                    # Transcribe - next 20% step
                                    if stepped_start != stepped_end:  # Only update if we move to a new step
                                        video_progress.progress(stepped_end)
                                    
                                    transcribe_msg.info(f"Transcribing video: {video_file.name}")
                                    result = model.transcribe(temp_path)
                                    transcribe_msg.empty()  # Clear the transcription message
                                    
                                    if "segments" not in result or not result["segments"]:
                                        st.error(f"No speech content found in video: {video_file.name}")
                                    else:
                                        transcript = ""
                                        metadata = []
                                        for segment in result["segments"]:
                                            start_time = segment["start"]
                                            text = segment["text"].strip()
                                            if text:
                                                transcript += f"{text} "
                                                metadata.append({"start_time": start_time})
                                        
                                        if transcript.strip():
                                            all_transcripts += transcript + "\n\n"
                                            all_video_metadata.extend(metadata)
                                            st.markdown(f'<div class="success-msg">✅ Successfully transcribed: {video_file.name}</div>', unsafe_allow_html=True)
                                        else:
                                            st.warning(f"No usable transcript from video: {video_file.name}")
                                except Exception as e:
                                    # Clear any remaining messages
                                    model_msg.empty()
                                    transcribe_msg.empty()
                                    st.error(f"Error transcribing video {video_file.name}: {str(e)}")
                        else:
                            # Single file processing - kept for backward compatibility
                            # Show stepped progress: 20%, 40%, 60%, 80%, 100%
                            video_progress.progress(20)  # Initial load
                        
                        # Ensure we end at 100%
                        video_progress.progress(100)
                        video_msg.empty()  # Clear the processing message
                        
                        if all_transcripts.strip():
                            video_chunks = get_text_chunks(all_transcripts)
                            all_chunks.extend(video_chunks)
                            all_metadata.extend(all_video_metadata[:len(video_chunks)])
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the video files.")

                # Process audio files
                if audio_files:
                    with audio_status.container():
                        # Create a placeholder for the processing message
                        audio_msg = st.empty()
                        audio_msg.write("⏳ Processing audio files...")
                        
                        # Add progress bar
                        audio_progress = st.progress(0)
                        all_transcripts = ""
                        all_audio_metadata = []
                        
                        if isinstance(audio_files, list):
                            total_audios = len(audio_files)
                            
                            for i, audio_file in enumerate(audio_files):
                                # Calculate step progress for this file
                                file_progress = int(((i+1) / total_audios) * 100)
                                # Round to nearest 20%
                                stepped_progress = ((file_progress + 10) // 20) * 20
                                audio_progress.progress(stepped_progress)
                                
                                # Create placeholder for model loading and transcription messages
                                model_msg = st.empty()
                                transcribe_msg = st.empty()
                                
                                try:
                                    # Model loading
                                    model_msg.info("Loading transcription model...")
                                    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_audio:
                                        temp_audio.write(audio_file.read())
                                        temp_path = temp_audio.name
                                    
                                    model = whisper.load_model("base")
                                    model_msg.empty()  # Clear the model loading message
                                    
                                    # Transcription
                                    transcribe_msg.info(f"Transcribing audio: {audio_file.name}")
                                    result = model.transcribe(temp_path)
                                    transcribe_msg.empty()  # Clear the transcription message
                                    
                                    if "segments" not in result or not result["segments"]:
                                        st.error(f"No speech content found in audio: {audio_file.name}")
                                    else:
                                        transcript = ""
                                        metadata = []
                                        for segment in result["segments"]:
                                            start_time = segment["start"]
                                            text = segment["text"].strip()
                                            if text:
                                                transcript += f"{text} "
                                                metadata.append({"start_time": start_time})
                                        
                                        if transcript.strip():
                                            all_transcripts += transcript + "\n\n"
                                            all_audio_metadata.extend(metadata)
                                            st.markdown(f'<div class="success-msg">✅ Successfully transcribed: {audio_file.name}</div>', unsafe_allow_html=True)
                                        else:
                                            st.warning(f"No usable transcript from audio: {audio_file.name}")
                                except Exception as e:
                                    # Clear any remaining messages
                                    model_msg.empty()
                                    transcribe_msg.empty()
                                    st.error(f"Error transcribing audio {audio_file.name}: {str(e)}")
                        else:
                            # Single file processing - show stepped progress
                            audio_progress.progress(20)  # Initial load
                            audio_progress.progress(40)  # Model loaded
                            audio_progress.progress(60)  # File processing
                            audio_progress.progress(80)  # Transcription
                        
                        # Ensure we end at 100%
                        audio_progress.progress(100)
                        audio_msg.empty()  # Clear the processing message
                        
                        if all_transcripts.strip():
                            audio_chunks = get_text_chunks(all_transcripts)
                            all_chunks.extend(audio_chunks)
                            all_metadata.extend(all_audio_metadata[:len(audio_chunks)])
                            files_processed = True
                        else:
                            st.warning("No text could be extracted from the audio files.")

                # Final processing
                if not files_processed:
                    st.warning("Please upload at least one file to process.")
                    return
                    
                if not all_chunks:
                    st.error("No content could be extracted from the uploaded files.")
                    return

                # Create vector store with a progress bar
                progress_container = st.container()
                with progress_container:
                    # st.markdown('<div style="background-color:#f0f9ff; padding:1rem; border-radius:8px; border-left:4px solid #0ea5e9;">', unsafe_allow_html=True)
                    st.write("🔍 Creating searchable index...")
                    progress_bar = st.progress(0)
                    
                    # Simulate progress for better UX
                    for percent_complete in range(100):
                        time.sleep(0.01)
                        progress_bar.progress(percent_complete + 1)
                    
                    get_vector_store(all_chunks, all_metadata)
                    
                    # Calculate total processing time
                    if st.session_state.process_start_time:
                        process_end_time = time.time()
                        total_duration = process_end_time - st.session_state.process_start_time
                        st.success("✅ Files processed and indexed! You can now ask questions about your documents.")
                        st.info(f"⏱️ Total processing time: {total_duration:.2f} seconds")
                        # Reset process start time
                        st.session_state.process_start_time = None
                    st.markdown('</div>', unsafe_allow_html=True)

            except Exception as e:
                st.error(f"⚠️ Error while processing files: {e}")
        
        # Add information at the bottom of the sidebar
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("""
        <hr style="margin: 0.3rem 0; padding: 0; border-color: #e2e8f0; opacity: 0.6;">
        """, unsafe_allow_html=True)
        st.markdown("""
        <div style="
            background: linear-gradient(90deg, #f0f9ff 0%, #e0f2fe 100%);
            padding: 1rem;
            border-radius: 8px;
            margin-top: 0.5rem;
            font-size: 0.9rem;
            color: #475569;
        ">
            <div style="text-align: center; margin-bottom: 0.5rem;">
                <span style="font-weight: 600;">Powered by</span>
            </div>
            <div style="display: flex; justify-content: space-around; align-items: center;">
                <div>Google Gemini</div>
                <div>OpenAI Whisper</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    # Add footer
    st.markdown(
        """
        <div class="footer">
            © 2024 Mediaverse - AI Powered Platform for Seamless Content Interaction | Developed  using Streamlit, Google Gemini, and OpenAI Whisper
        </div>
        """,
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()